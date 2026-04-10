from pptx import Presentation
from docx import Document

# Create PPTX
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'AI Office Assistant'
slide.placeholders[1].text = 'Offline AI powered document processing system'
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = 'Features'
slide2.placeholders[1].text = 'PDF, DOCX, PPTX, TXT summarization using T5 and Mistral models'
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
slide3.shapes.title.text = 'Tech Stack'
slide3.placeholders[1].text = 'FastAPI backend, Flutter frontend, T5-small model, Mistral 7B via Ollama'
prs.save('test_ppt.pptx')
print('test_ppt.pptx created!')

# Create DOCX
doc = Document()
doc.add_heading('AI Office Assistant — Test Document', 0)
doc.add_heading('Project Overview', level=1)
doc.add_paragraph('This is a test Word document for the AI Office Assistant benchmark suite.')
doc.add_paragraph('The system processes DOCX files by extracting text from paragraphs and tables.')
doc.add_heading('Technical Details', level=1)
doc.add_paragraph('Dynamic chunking is used based on document size to optimize performance.')
doc.add_paragraph('T5-small model handles summarization for fast offline processing.')
doc.add_paragraph('Mistral 7B via Ollama provides higher quality structured summaries.')
doc.add_heading('Supported Formats', level=1)
doc.add_paragraph('PDF, DOCX, DOC, PPTX, PPT, TXT files are all supported.')
doc.save('test_doc.docx')
print('test_doc.docx created!')

print('All test files created successfully!')
