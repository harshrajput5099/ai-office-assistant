# backend/utils/document_extractor.py
"""
Unified Document Extractor
Supports: PDF, DOCX, DOC, PPTX, PPT, TXT
Returns plain text from any supported format.
"""

import io
import os
import time
import logging

logger = logging.getLogger(__name__)

# ─── SUPPORTED FORMATS ────────────────────────────────────
SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.doc', '.pptx', '.ppt', '.txt'}


# ══════════════════════════════════════════════════════════
# MAIN EXTRACTOR — routes to correct handler
# ══════════════════════════════════════════════════════════
def extract_text(file_bytes: bytes, filename: str) -> dict:
    """
    Master extractor. Detects file type and routes to correct handler.
    Returns dict with: text, page_count, file_size_kb, extraction_time_s
    """
    start = time.time()
    ext = os.path.splitext(filename.lower())[1]

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type: '{ext}'. "
            f"Supported: {', '.join(SUPPORTED_EXTENSIONS)}"
        )

    file_size_kb = round(len(file_bytes) / 1024, 2)

    if ext == '.pdf':
        text, page_count = _extract_pdf(file_bytes)
    elif ext in ('.docx', '.doc'):
        text, page_count = _extract_docx(file_bytes, filename)
    elif ext in ('.pptx', '.ppt'):
        text, page_count = _extract_pptx(file_bytes)
    elif ext == '.txt':
        text, page_count = _extract_txt(file_bytes)
    else:
        raise ValueError(f"Handler missing for: {ext}")

    extraction_time = round(time.time() - start, 3)

    if not text or len(text.strip()) < 10:
        raise ValueError(
            f"No readable text found in '{filename}'. "
            "File may be scanned/image-based or empty."
        )

    logger.info(
        f"Extracted '{filename}' | "
        f"{page_count} pages | {file_size_kb}KB | {extraction_time}s"
    )

    return {
        "text": text,
        "page_count": page_count,
        "file_size_kb": file_size_kb,
        "extraction_time_s": extraction_time,
        "filename": filename,
        "extension": ext,
    }


# ══════════════════════════════════════════════════════════
# PDF EXTRACTOR
# ══════════════════════════════════════════════════════════
def _extract_pdf(file_bytes: bytes) -> tuple[str, int]:
    try:
        import pdfplumber
    except ImportError:
        raise ImportError("Run: pip install pdfplumber")

    text_parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        page_count = len(pdf.pages)
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_parts.append(page_text)

    return '\n\n'.join(text_parts), page_count


# ══════════════════════════════════════════════════════════
# DOCX EXTRACTOR
# ══════════════════════════════════════════════════════════
def _extract_docx(file_bytes: bytes, filename: str) -> tuple[str, int]:
    ext = os.path.splitext(filename.lower())[1]

    # .doc files — convert to docx first using python-docx2txt
    if ext == '.doc':
        return _extract_doc_legacy(file_bytes)

    try:
        from docx import Document
    except ImportError:
        raise ImportError("Run: pip install python-docx")

    doc = Document(io.BytesIO(file_bytes))
    paragraphs = []

    # Extract from body paragraphs
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)

    # Extract from tables
    for table in doc.tables:
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_texts:
                paragraphs.append(' | '.join(row_texts))

    # Estimate page count (approx 500 words per page)
    total_words = sum(len(p.split()) for p in paragraphs)
    estimated_pages = max(1, total_words // 500)

    return '\n\n'.join(paragraphs), estimated_pages


def _extract_doc_legacy(file_bytes: bytes) -> tuple[str, int]:
    """Extract from old .doc format using docx2txt"""
    try:
        import docx2txt
        import tempfile
    except ImportError:
        raise ImportError("Run: pip install docx2txt")

    with tempfile.NamedTemporaryFile(delete=False, suffix='.doc') as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        text = docx2txt.process(tmp_path)
    finally:
        os.unlink(tmp_path)

    total_words = len(text.split())
    estimated_pages = max(1, total_words // 500)
    return text, estimated_pages


# ══════════════════════════════════════════════════════════
# PPTX EXTRACTOR
# ══════════════════════════════════════════════════════════
def _extract_pptx(file_bytes: bytes) -> tuple[str, int]:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("Run: pip install python-pptx")

    prs = Presentation(io.BytesIO(file_bytes))
    slide_texts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            # Extract from tables inside slides
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = ' | '.join(
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text.strip()
                    )
                    if row_text:
                        texts.append(row_text)

        if texts:
            slide_texts.append(f"[Slide {slide_num}]\n" + '\n'.join(texts))

    page_count = len(prs.slides)
    return '\n\n'.join(slide_texts), page_count


# ══════════════════════════════════════════════════════════
# TXT EXTRACTOR
# ══════════════════════════════════════════════════════════
def _extract_txt(file_bytes: bytes) -> tuple[str, int]:
    # Try UTF-8 first, fallback to latin-1
    try:
        text = file_bytes.decode('utf-8')
    except UnicodeDecodeError:
        text = file_bytes.decode('latin-1', errors='replace')

    total_words = len(text.split())
    estimated_pages = max(1, total_words // 500)
    return text, estimated_pages